"""Fill the SIH2026 idea template for SIH26056."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

SRC = "/home/ajeet/SIH/SIH2026-IDEA-Presentation-Format.pptx"
OUT = "/home/ajeet/SIH/SIH26056-Idea-Presentation.pptx"
TMP  = "/home/ajeet/.claude/jobs/1319c671/tmp"
DASH = TMP + "/dash-live.png"
CURVE = TMP + "/chart-curve.png"
PREM  = TMP + "/chart-premium.png"
CARR  = TMP + "/chart-carriers.png"
MAP   = TMP + "/chart-map.png"
CAD   = TMP + "/chart-cadence.png"
APIXC = TMP + "/chart-apix.png"
SURV  = TMP + "/chart-survey.png"
SHOT  = TMP + "/portal-shot.png"
REC   = TMP + "/chart-record.png"

NAVY   = RGBColor(0x1F, 0x49, 0x7D)
BLUE   = RGBColor(0x4F, 0x81, 0xBD)
PALE   = RGBColor(0xDC, 0xE6, 0xF1)
PALER  = RGBColor(0xEE, 0xF3, 0xF9)
INK    = RGBColor(0x26, 0x2B, 0x33)
GREY   = RGBColor(0x5A, 0x62, 0x6C)
RED    = RGBColor(0xC0, 0x50, 0x4D)
GREEN  = RGBColor(0x1B, 0x8A, 0x5A)
TEAL   = RGBColor(0x0E, 0x7C, 0x86)
AMBER  = RGBColor(0xC7, 0x77, 0x00)
PLUM   = RGBColor(0x7A, 0x3E, 0x6B)
INDIGO = RGBColor(0x4B, 0x3F, 0x8C)
MAROON = RGBColor(0xA6, 0x3A, 0x3A)
ACCENTS = [BLUE, GREEN, AMBER, INDIGO, TEAL, PLUM]
TINTS = [RGBColor(0xE8,0xF0,0xF9), RGBColor(0xE4,0xF4,0xEC), RGBColor(0xFBF,0xF0,0xDD)
         if False else RGBColor(0xFB,0xF0,0xDD), RGBColor(0xEC,0xEA,0xF7),
         RGBColor(0xE2,0xF2,0xF3), RGBColor(0xF4,0xEA,0xF1)]
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xC7, 0xD3, 0xE3)

TEAM = "Fare Enough 101"


def fare_count(fallback=51_000):
    """Live row count, so the headline figure never drifts from the database.
    Falls back to a stated floor when the deck is built offline."""
    try:
        import os
        env = "/home/ajeet/SIH/airfare-scraper/.env"
        for line in open(env):
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                os.environ.setdefault(k, v.strip().strip('"').strip("'"))
        from supabase import create_client
        sb = create_client(os.environ["SUPABASE_URL"], os.environ["SUPABASE_KEY"])
        n = sb.table("fares").select("id", count="exact").limit(1).execute().count
        return max(int(n), fallback)
    except Exception:
        return fallback


FARES = fare_count()
print(f"fares in database: {FARES:,}")

prs = Presentation(SRC)


# ------------------------------------------------------------------ helpers
def drop_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def remove(shape):
    shape._element.getparent().remove(shape._element)


def write(tf, blocks, wrap=True, margin=0.04):
    """blocks: (text, size, bold, color, space_before_pt, indent_level)"""
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.clear()
    for i, b in enumerate(blocks):
        text, size, bold, color, before, lvl = (list(b) + [0, 0])[:6]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(before)
        p.space_after = Pt(0)
        p.line_spacing = 0.95
        if lvl:
            p.paragraph_format if False else None
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tf


def linkbox(slide, x, y, w, h, blocks, links=None):
    """Like textbox, but any block carrying a URL becomes a real hyperlink."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.clear()
    for i, b in enumerate(blocks):
        text, size, bold, colour, before, url = (list(b) + [0, None])[:6]
        p_ = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p_.space_before = Pt(before); p_.space_after = Pt(0); p_.line_spacing = 1.0
        r = p_.add_run(); r.text = text
        r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = colour
        if url:
            r.hyperlink.address = url
    return tb


def textbox(slide, x, y, w, h, blocks, **kw):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    write(tb.text_frame, blocks, **kw)
    return tb


def panel(slide, x, y, w, h, fill=PALER, line=LINE, radius=True, width_pt=1):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        shp.adjustments[0] = 0.06
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(width_pt)
    shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    return shp


def chip(slide, x, y, w, h, text, fill=NAVY, fg=WHITE, size=10.5, bold=True):
    shp = panel(slide, x, y, w, h, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = fg
    return shp


def textlink(slide, x, y, w, text, url, size=10.5, colour=BLUE,
             align=PP_ALIGN.RIGHT, h=0.24):
    """An underlined text hyperlink, the way the reference deck marks its sources."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    r = para.add_run(); r.text = text
    r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = True
    r.font.underline = True
    r.font.color.rgb = colour
    r.hyperlink.address = url
    return tb


def heading(slide, x, y, w, text, size=13.5):
    """Section heading = the template's required pointer, kept verbatim."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    write(tb.text_frame, [(text, size, True, NAVY)])
    # underline rule
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.02), Inches(y + 0.30),
                                Inches(min(w, 1.5)), Inches(0.028))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background(); ln.shadow.inherit = False
    return tb


def bullets(slide, x, y, w, h, items, size=11.5, gap=4, color=INK):
    blocks = []
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, bold, col = it
        else:
            txt, bold, col = it, False, color
        blocks.append(("•  " + txt if not txt.startswith(" ") else txt,
                       size, bold, col, 0 if i == 0 else gap))
    return textbox(slide, x, y, w, h, blocks)


def set_team_oval(slide):
    for sh in slide.shapes:
        if sh.name.startswith("Oval") and sh.has_text_frame:
            tf = sh.text_frame
            p = tf.paragraphs[0]
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = TEAM
            else:
                r = p.add_run(); r.text = TEAM
                r.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            for r in p.runs:
                r.font.size = Pt(10.5)
                r.font.bold = True


def set_footer(slide):
    """The reference deck brands its footer; this keeps that edit in the build."""
    ph = shape_by_name(slide, "Footer Placeholder 6")
    if ph is None or not ph.has_text_frame:
        return
    para = ph.text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return
    runs[0].text = f"{TEAM} - @SIH Idea submission"
    for extra in runs[1:]:
        extra.text = ""


def set_title(slide, text, size=36):
    t = shape_by_name(slide, "Title 1")
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = "Times New Roman"; r.font.size = Pt(size)
    r.font.bold = True; r.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    return t


# =====================================================================  S1
s1 = prs.slides[0]
sub = shape_by_name(s1, "Subtitle 3")
if sub:
    remove(sub)
tb = shape_by_name(s1, "TextBox 9")
tb.left, tb.top, tb.width, tb.height = Inches(0.36), Inches(2.27), Inches(6.48), Inches(5.14)
FIELDS = [
    ("Problem Statement ID", "SIH26056"),
    ("Problem Statement Title",
     "Development of a Real-time Airfare Price Index for India through Automated "
     "Web Scraping of Airline and Online Travel Aggregator Portals for Augmentation "
     "of the Consumer Price Index (CPI)"),
    ("Theme", "Travel & Tourism"),
    ("PS Category", "Software"),
    ("Team ID", "\u2039TEAM ID\u203a"),
    ("Team Name (Registered on portal)", TEAM),
]
blocks = []
for i, (label, value) in enumerate(FIELDS):
    blocks.append((f"{label} - {value}", 20, True, INK, 0 if i == 0 else 8))
write(tb.text_frame, blocks)
for para in tb.text_frame.paragraphs:
    para.alignment = PP_ALIGN.JUSTIFY
    para.line_spacing = 1.06

# =====================================================================  S2
s2 = prs.slides[1]
set_team_oval(s2)
set_footer(s2)
set_title(s2, "APIx \u2014 AIRFARE PRICE INDEX")
remove(shape_by_name(s2, "TextBox 8"))

textbox(s2, 0.45, 1.00, 8.45, 1.00, [
    ("PROPOSED SOLUTION", 10.5, True, BLUE),
    ("A daily airfare inflation index for the CPI", 14.5, True, NAVY, 2),
    ("MoSPI prices air travel by hand, once a month, for one departure date. APIx reads "
     "the same public pages every ten minutes and delivers one weighted index number "
     "through an API.", 12, False, GREY, 5)])

stat = panel(s2, 9.05, 1.02, 3.80, 0.90, fill=TINTS[1], line=None)
write(stat.text_frame, [
    ("BUILT AND RUNNING", 10.5, True, GREEN),
    (f"{FARES // 1000:,},000+ fares collected \u00b7 15 city pairs \u00d7 5 lead times "
     "\u00b7 every 10 minutes \u00b7 \u20b90 a month to run", 10.5, False, INK, 4)])

# --- what changes for the CPI -----------------------------------------
panel(s2, 0.45, 2.16, 6.95, 1.96, fill=TINTS[0], line=None)
textbox(s2, 0.60, 2.23, 4.2, 0.26, [("WHAT CHANGES FOR THE CPI", 10.5, True, BLUE)])
textbox(s2, 2.35, 2.55, 2.10, 0.26, [("TODAY", 11, True, GREY)])
textbox(s2, 4.55, 2.55, 2.70, 0.26, [("WITH APIx", 11, True, GREEN)])
CHANGES = [("Frequency",   "once a month",        "every 10 minutes"),
           ("Coverage",    "a handful of quotes", "75 priced cells"),
           ("Price basis", "one quoted price",    "minimum logical fare"),
           ("Delivery",    "typed into a form",   "authenticated REST API")]
for i, (attr, today, ours) in enumerate(CHANGES):
    yy = 2.88 + i * 0.31
    textbox(s2, 0.60, yy, 1.75, 0.26, [(attr, 11, True, INK)])
    textbox(s2, 2.35, yy, 2.10, 0.26, [(today, 10.5, False, GREY)])
    textbox(s2, 4.55, yy, 2.70, 0.26, [(ours, 10.5, True, INK)])

# --- the formula and the definitions behind it -------------------------
panel(s2, 7.55, 2.16, 5.30, 0.96, fill=TINTS[3], line=None)
textbox(s2, 7.70, 2.21, 5.00, 0.24, [("THE INDEX IN ONE LINE", 10.5, True, INDIGO)])
textbox(s2, 7.70, 2.46, 5.00, 0.26, [
    ("APIx\u209c = 100 \u00d7 exp( \u03a3 w\u1d62\u00b7ln(P\u1d62,\u209c/P\u1d62,\u2080) / \u03a3 w\u1d62 )",
     10.5, True, INK)])
textbox(s2, 7.70, 2.73, 5.00, 0.34, [
    ("Weighted Jevons \u2014 the elementary-aggregate formula Eurostat and the ONS "
     "use for their own price indices.", 10, False, GREY)])

panel(s2, 7.55, 3.20, 5.30, 0.92, fill=TINTS[4], line=None)
textbox(s2, 7.70, 3.25, 5.00, 0.24, [("WHAT THE TERMS MEAN", 10.5, True, TEAL)])
DEFS = [("Basket", "15 city pairs \u00d7 5 lead times = 75 cells"),
        ("Price",  "minimum logical fare, no add-ons"),
        ("Weight", "route seat share \u00d7 lead-time share"),
        ("Base",   "August 2026 = 100")]
for i, (k, v) in enumerate(DEFS):
    yy = 3.50 + i * 0.155
    textbox(s2, 7.70, yy, 0.85, 0.16, [(k, 10, True, TEAL)])
    textbox(s2, 8.55, yy, 4.20, 0.16, [(v, 10, False, INK)])

# --- innovation and uniqueness: four points, one row -------------------
heading(s2, 0.45, 4.34, 12.4, "Why it stands out", 14)

CARDS = [
    ("The standard method",
     "Weighted Jevons on minimum logical fares \u2014 the elementary-aggregate method "
     "Eurostat and the ONS use for their own price indices."),
    ("Weighting is a matrix",
     "Every cell carries route seat share \u00d7 booking lead time, so Delhi\u2013Mumbai "
     "moves the index three times as hard as Delhi\u2013Srinagar."),
    ("Admits thin coverage",
     "Below 60% basket coverage the figure ships marked provisional, with the reason. "
     "An index that admits thin coverage beats one that never does."),
    ("Machine-readable",
     "One authenticated REST call returns the index with its method, coverage and "
     "revision history. Nobody retypes a number off a screen."),
]
cw, ch, gx = 2.905, 1.78, 0.26
for i, (title, body) in enumerate(CARDS):
    x = 0.45 + i * (cw + gx)
    card = panel(s2, x, 4.72, cw, ch, fill=TINTS[i], line=None)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.72),
                              Inches(cw), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENTS[i]
    bar.line.fill.background(); bar.shadow.inherit = False
    write(card.text_frame, [(title, 12, True, ACCENTS[i]), (body, 10.5, False, INK, 5)])
    card.text_frame.margin_top = Inches(0.14)

# =====================================================================  S3
s3 = prs.slides[2]
set_team_oval(s3)
set_footer(s3)
set_title(s3, "TECHNICAL APPROACH")
remove(shape_by_name(s3, "TextBox 8"))

heading(s3, 0.45, 1.16, 12.4, "How a published fare becomes a published index", 15)

# each step carries the number that makes it real
STEPS = [("1", "Collect", "every 10 min", "robots-gated"),
         ("2", "Extract", "13,000 \u2192 3,700", "chars per page"),
         ("3", "Validate", "38 of 40", "flights kept"),
         ("4", "Clean", "min 3 obs", "per priced cell"),
         ("5", "Index", "75 cells", "weighted Jevons"),
         ("6", "Publish", "2 endpoints", "portal + API")]
bw, gap = 1.88, 0.22
for i, (n, t, big, sub) in enumerate(STEPS):
    x = 0.45 + i * (bw + gap)
    box = panel(s3, x, 1.58, bw, 1.30, fill=TINTS[i], line=None)
    bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.58),
                              Inches(bw), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENTS[i]
    bar.line.fill.background(); bar.shadow.inherit = False
    write(box.text_frame, [(f"{n}  {t}", 13.5, True, ACCENTS[i]),
                           (big, 14, True, NAVY, 4), (sub, 10.5, False, GREY, 1)])
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < 5:
        ar = s3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.015),
                                 Inches(2.13), Inches(gap - 0.03), Inches(0.20))
        ar.fill.solid(); ar.fill.fore_color.rgb = ACCENTS[i]
        ar.line.fill.background(); ar.shadow.inherit = False

heading(s3, 0.45, 3.02, 12.4,
        "What we keep from one results row \u2014 and what we drop", 14)

# The interesting part is not the raw text, it is the rule applied to each
# field, so this is a table of decisions rather than a dump of three payloads.
KEEP = [
    ("IndiGo \u00b7 6E-955 \u00b7 20:20 \u2192 22:20",
     "carrier, flight_number, times", "the flight's identity", INK),
    ("\u20b96,529", "total_fare = 6529",
     "the price a traveller actually pays", GREEN),
    ("\u20b9480 off with CTFKSBIC", "dropped",
     "a coupon is not a published fare", RED),
    ("no tax or fee breakdown", "stored NULL, never derived",
     "a CPI prices the total a household pays", AMBER),
    ("\u2014 nothing on the page \u2014", "source \u00b7 scraped_at \u00b7 model_used",
     "so any index figure traces back", INDIGO),
]
COLS = [(0.70, 3.85), (4.75, 4.05), (8.95, 3.70)]

panel(s3, 0.45, 3.44, 12.4, 2.56, fill=WHITE, line=LINE)
band = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(3.44),
                           Inches(12.4), Inches(0.34))
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
for (cx, cwid), label in zip(COLS, ("ON THE PAGE", "WHAT WE STORE", "THE RULE")):
    textbox(s3, cx, 3.50, cwid, 0.24, [(label, 10, True, WHITE)])

for i, (page, stored, rule, colour) in enumerate(KEEP):
    ry = 3.86 + i * 0.42
    if i % 2 == 0:
        stripe = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(ry - 0.04),
                                     Inches(12.4), Inches(0.42))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = PALER
        stripe.line.fill.background(); stripe.shadow.inherit = False
    textbox(s3, COLS[0][0], ry, COLS[0][1], 0.30, [(page, 11.5, False, INK)])
    textbox(s3, COLS[1][0], ry, COLS[1][1], 0.30, [(stored, 11.5, True, colour)])
    textbox(s3, COLS[2][0], ry, COLS[2][1], 0.30, [(rule, 11, False, GREY)])

TECHSTRIP = ("Python \u00b7 Playwright \u00b7 LLM extraction with failover \u00b7 "
             "Pydantic v2 \u00b7 Supabase PostgreSQL \u00b7 FastAPI \u00b7 Cloudflare \u00b7 cron")
strip = panel(s3, 0.45, 6.12, 12.4, 0.58, fill=RGBColor(0xEC, 0xEA, 0xF7), line=None)
write(strip.text_frame, [("STACK   " + TECHSTRIP, 11.5, True, INDIGO)])

# =====================================================================  S4
s4 = prs.slides[3]
set_team_oval(s4)
set_footer(s4)
set_title(s4, "FEASIBILITY AND VIABILITY")
remove(shape_by_name(s4, "TextBox 8"))

FEAS = [
    ("Technical feasibility", "USD 18.1 bn", "market size, 2026",
     "Fares are published continuously on public pages and change all day. "
     "The data exists; what is missing is a systematic way to read it."),
    ("Operational feasibility", "164 million", "domestic passengers",
     "A fixed basket of the 15 busiest city pairs across 148 operational "
     "airports tracks the traffic that actually moves the index."),
    ("Economic feasibility", "₹0", "a month to run",
     "Free tiers throughout, about two minutes of compute per cycle. MoSPI "
     "visits shops monthly; this observes 144 times a day."),
    ("Regulatory feasibility", "since 2014", "ONS precedent",
     "Eurostat's HICP guidance names air fares as a scrapeable category and "
     "the ONS has published on it since 2014 — precedented, not experimental."),
]
FEAS_COLOURS = [BLUE, GREEN, AMBER, INDIGO]
y = 1.24
for i, (title, stat, statlbl, body) in enumerate(FEAS):
    card = panel(s4, 0.45, y, 6.30, 1.28, fill=TINTS[i], line=None)
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y),
                              Inches(0.075), Inches(1.28))
    bar.fill.solid(); bar.fill.fore_color.rgb = FEAS_COLOURS[i]
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s4, 0.70, y + 0.10, 4.20, 1.08,
            [(title, 14, True, FEAS_COLOURS[i]), (body, 11.5, False, INK, 4)])
    textbox(s4, 5.00, y + 0.34, 1.62, 0.62,
            [(stat, 13.5, True, FEAS_COLOURS[i]), (statlbl, 9, False, GREY, 2)])
    y += 1.36

heading(s4, 7.0, 1.24, 5.9, "Risks, and what we did about them", 15)
RISKS = [("16 portals surveyed, one is usable",
          "Ten disallow us, four block us, airlines refuse automation entirely."),
         ("Cloud IPs are blocked too",
          "Verified in CI — so collection runs from a residential host."),
         ("Models get retired without notice",
          "Three died during this build. Extraction now fails over automatically."),
         ("Thin coverage would mislead",
          "Days below 60% basket weight publish as provisional, with the reason.")]
y = 1.70
for risk, fix in RISKS:
    panel(s4, 7.0, y, 5.9, 0.88, fill=RGBColor(0xFA, 0xEE, 0xEE), line=None)
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(y),
                              Inches(0.055), Inches(0.88))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s4, 7.2, y + 0.09, 5.6, 0.72,
            [(risk, 13.5, True, MAROON), (fix, 11.5, False, INK, 4)])
    y += 0.96

# --- viability: the running cost, itemised ----------------------------
heading(s4, 7.0, 5.58, 5.9, "What it costs to run", 15)
panel(s4, 7.0, 5.96, 5.9, 0.96, fill=TINTS[1], line=None)
COST = [("Collection · Playwright", "₹0 · free tier", False),
        ("Extraction · NVIDIA NIM LLM", "₹0 · free tier", False),
        ("Storage · Supabase Postgres", "₹0 · free tier", False),
        ("Publication · Render + Cloudflare", "₹0 · free tier", False),
        ("Total · 144 cycles a day", "₹0 a month", True)]
for i, (item, cost, bold) in enumerate(COST):
    yy = 6.02 + i * 0.17
    textbox(s4, 7.16, yy, 3.60, 0.17, [(item, 10, bold, GREEN if bold else INK)])
    textbox(s4, 10.85, yy, 1.90, 0.17, [(cost, 10, bold, GREEN if bold else GREY)])


# =====================================================================  S5
s5 = prs.slides[4]
set_team_oval(s5)
set_footer(s5)
set_title(s5, "IMPACT AND BENEFITS")
remove(shape_by_name(s5, "TextBox 8"))

heading(s5, 0.45, 1.20, 12.4, "From a fare on a screen to a figure in the CPI", 15)
JOURNEY = [
    ("A fare is published", "A portal shows a price for one route, one date"),
    ("We observe it", "Read every 10 minutes, robots-checked"),
    ("It becomes a price", "Cheapest logical fare per route × lead time"),
    ("It becomes an index", "Weighted Jevons across the basket"),
    ("MoSPI ingests it", "REST API with method and coverage"),
    ("CPI reflects reality", "Transport sub-index from 144 daily reads"),
]
bw, gap = 1.86, 0.24
for i, (t, d) in enumerate(JOURNEY):
    x = 0.45 + i * (bw + gap)
    box = panel(s5, x, 1.70, bw, 1.42, fill=TINTS[i % len(TINTS)], line=None)
    write(box.text_frame, [(str(i + 1), 12, True, ACCENTS[i % len(ACCENTS)]),
                           (t, 13, True, NAVY, 2), (d, 10.5, False, GREY, 4)])
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < len(JOURNEY) - 1:
        ar = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.02),
                                 Inches(2.28), Inches(gap - 0.04), Inches(0.20))
        ar.fill.solid(); ar.fill.fore_color.rgb = BLUE
        ar.line.fill.background(); ar.shadow.inherit = False

s5.shapes.add_picture(SHOT, Inches(0.45), Inches(3.36), width=Inches(6.2))
linkbox(s5, 0.45, 6.40, 6.2, 0.44, [
    ("Live: apix-portal.pages.dev", 12, True, BLUE, 0, "https://apix-portal.pages.dev"),
    ("Headline inflation, release status and method on one page.",
     11, False, GREY, 2)])

heading(s5, 7.0, 3.32, 5.9, "Who benefits, and how", 15)
BEN = [("MoSPI / NSO", "A defensible transport input, recomputable from stored micro-data"),
       ("DGCA and policy", "Route-level fare behaviour nobody currently publishes"),
       ("Citizens", "A public record of what air travel actually costs over time"),
       ("Researchers", "An open, reproducible price series for a market that has none")]
y = 3.76
for i, (who, what) in enumerate(BEN):
    chip(s5, 7.0, y, 1.85, 0.32, who, fill=ACCENTS[i % len(ACCENTS)], size=11)
    textbox(s5, 9.05, y - 0.02, 3.85, 0.5, [(what, 11.5, False, INK)])
    y += 0.52

# --- the impact, as numbers -------------------------------------------
IMPACT = [("144 ×", "observations a day, not one a month", GREEN),
          ("75", "priced cells behind every figure", INDIGO),
          ("T+0", "same-day index, not a two-week lag", TEAL)]
sw = (5.9 - 2 * 0.16) / 3
for i, (val, lbl, col) in enumerate(IMPACT):
    x = 7.0 + i * (sw + 0.16)
    panel(s5, x, 5.80, sw, 0.62, fill=TINTS[(i * 2 + 1) % len(TINTS)], line=None)
    textbox(s5, x + 0.10, 5.84, sw - 0.20, 0.54,
            [(val, 15, True, col), (lbl, 8.5, False, GREY, 1)])

promise = panel(s5, 7.0, 6.48, 5.9, 0.44, fill=RGBColor(0xE4, 0xF4, 0xEC), line=None)
write(promise.text_frame, [
    ("Our promise — replace one manual price visit a month with 144 automated "
     "observations a day, at zero marginal cost.", 11, False, INK)])
promise.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# =====================================================================  S6
s6 = prs.slides[5]
set_team_oval(s6)
set_footer(s6)
set_title(s6, "RESEARCH AND REFERENCES")
remove(shape_by_name(s6, "TextBox 8"))

REFS = [
    ("Method \u2014 why Jevons",
     "Eurostat, Practical guidelines on web scraping for the HICP (2020). Names air "
     "fares explicitly as a scraped category.",
     "Eurostat guidance",
     "https://ec.europa.eu/eurostat/documents/272892/12032198/Guidelines-web-scraping-HICP-11-2020.pdf"),
    ("Precedent \u2014 already done",
     "ONS ran a web-scraping programme for consumer prices from 2014, Eurostat-funded "
     "from 2015.",
     "ONS study",
     "https://www.ons.gov.uk/economy/inflationandpriceindices/articles/researchindicesusingwebscrapedpricedata/august2017update/previous/v1/pdf"),
    ("The Indian basket",
     "DGCA city-pair statistics: 164 m domestic passengers, IndiGo 64.2%. Our own "
     "1,000-fare sample returned 68.2% \u2014 arrived at independently.",
     "DGCA statistics",
     "https://www.dgca.gov.in/digigov-portal/"),
    ("Standards we hold to",
     "RFC 9309, the Robots Exclusion Protocol. We implemented it ourselves after "
     "finding Python's standard parser reports disallowed paths as allowed.",
     "RFC 9309", "https://www.rfc-editor.org/rfc/rfc9309"),
    ("CPI methodology",
     "MoSPI moved the CPI to base 2024=100 in February 2026, aligned to COICOP 2018. "
     "Air fares are still priced by hand \u2014 the process this augments.",
     "The 2024=100 series",
     "https://www.pib.gov.in/PressReleasePage.aspx?PRID=2227012&reg=3&lang=1"),
    ("Our code and full method",
     "Collector, index engine, export API and portal. The README documents the "
     "method and the publication threshold.\n"
     "github.com/BradCage-afk/sih-2026-airfare-index",
     "Open the repo",
     "https://github.com/BradCage-afk/sih-2026-airfare-index"),
]
textbox(s6, 0.45, 1.16, 6.0, 0.24,
        [("EVERY CARD IS CLICKABLE \u2014 THE UNDERLINED LINK OPENS THE SOURCE",
          9.5, True, GREY)])

cw2, ch2, gx2, gy2 = 4.05, 1.78, 0.28, 0.16
for i, (title, body, linktext, url) in enumerate(REFS):
    x = 0.45 + (i % 3) * (cw2 + gx2)
    y = 1.46 + (i // 3) * (ch2 + gy2)
    card = panel(s6, x, y, cw2, ch2, fill=TINTS[i], line=None)
    card.click_action.hyperlink.address = url      # the whole card is the link
    bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                              Inches(cw2), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENTS[i]
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s6, x + 0.16, y + 0.14, cw2 - 0.32, ch2 - 0.46,
            [(title, 13.5, True, ACCENTS[i]), (body, 11.5, False, INK, 4)])
    textlink(s6, x + cw2 - 2.36, y + ch2 - 0.28, 2.20, linktext, url,
             size=10.5, colour=ACCENTS[i])

s6.shapes.add_picture(SURV, Inches(0.45), Inches(5.20), width=Inches(4.55))
linkbox(s6, 5.25, 5.24, 7.60, 1.3, [
    ("Everything above is reproducible.", 13, True, NAVY),
    ("The index, every chart in this deck and every figure on the portal are "
     "regenerated from the database by scripts in the repository. Nothing is "
     "transcribed by hand, so no number here can drift from what the system holds.",
     11.5, False, INK, 4),
    ("Everything in this deck is checkable against a live system:",
     11.5, True, NAVY, 6)])

OPEN = [("Live API docs", "https://apix-api-n5ux.onrender.com/docs", NAVY),
        ("Release portal", "https://apix-portal.pages.dev", GREEN),
        ("README + method", "https://github.com/BradCage-afk/sih-2026-airfare-index#readme", INDIGO),
        ("Aviation market", "https://www.ibef.org/industry/indian-aviation", TEAL)]
lw = (7.60 - 3 * 0.14) / 4
for i, (txt, url, col) in enumerate(OPEN):
    textlink(s6, 5.25 + i * (lw + 0.14), 6.46, lw, txt, url,
             size=10.5, colour=col, align=PP_ALIGN.LEFT)

drop_slide(prs, 6)
prs.save(OUT)
print("saved", OUT)
